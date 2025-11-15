from __future__ import annotations

import os
from pathlib import Path
from typing import List, Dict
import re


def _html_to_text(html: str) -> str:
    try:
        from bs4 import BeautifulSoup  # type: ignore
    except Exception:
        # Fallback: naive tag strip
        import re

        return re.sub(r"<[^>]+>", " ", html or "").strip()
    soup = BeautifulSoup(html or "", "html.parser")
    return soup.get_text("\n", strip=True)


def _clean_extracted_text(text: str) -> str:
    """
    Clean common PDF extraction artifacts to improve chunk quality.
    """
    if not text:
        return ""
    # Remove <latexit ...>...</latexit> blocks (often embedded LaTeX artifacts)
    text = re.sub(r"<\s*latexit[^>]*>.*?<\s*/\s*latexit\s*>", " ", text, flags=re.DOTALL | re.IGNORECASE)
    # Remove long base64-like blobs
    text = re.sub(r"[A-Za-z0-9+/=]{80,}", " ", text)
    # Normalize ligatures
    text = text.replace("\ufb01", "fi").replace("\ufb02", "fl")
    # Remove zero-width characters and BOM
    text = re.sub(r"[\u200b\u200c\u200d\ufeff]", "", text)
    # Fix hyphenation across line breaks: "exam-\nple" -> "example"
    text = re.sub(r"-\s*\n\s*", "", text)
    # Normalize whitespace
    text = re.sub(r"[ \t]+", " ", text)
    # Trim trailing spaces per line
    text = "\n".join(ln.rstrip() for ln in text.splitlines())
    return text


def _pdf_to_text(path: str) -> str:
    # Primary: PyPDF2
    try:
        from PyPDF2 import PdfReader  # type: ignore
        text = ""
        with open(path, "rb") as f:
            reader = PdfReader(f)
            for idx, page in enumerate(reader.pages, start=1):
                text += f"=== Page {idx} ===\n"
                text += (page.extract_text() or "") + "\n"
        if text.strip():
            return _clean_extracted_text(text)
    except Exception:
        pass
    # Fallback: PyMuPDF (fitz) if available
    try:
        import fitz  # type: ignore
        text = ""
        with fitz.open(path) as doc:
            for idx, p in enumerate(doc, start=1):
                text += f"=== Page {idx} ===\n"
                text += (p.get_text() or "") + "\n"
        if text.strip():
            return _clean_extracted_text(text)
    except Exception:
        pass
    # Fallback: pdfminer.six if available
    try:
        from pdfminer.high_level import extract_text  # type: ignore
        text = extract_text(path) or ""
        if text.strip():
            return _clean_extracted_text(text)
    except Exception:
        pass
    return ""


def _docx_to_text(path: str) -> str:
    try:
        import docx  # type: ignore
    except Exception:
        return ""
    try:
        d = docx.Document(path)
        lines = []
        for para in d.paragraphs:
            txt = para.text or ""
            if not txt.strip():
                continue
            # Heading styles: "Heading 1".."Heading 6"
            try:
                style_name = getattr(para.style, "name", "") or ""
            except Exception:
                style_name = ""
            level = None
            if style_name.startswith("Heading"):
                try:
                    n = int(style_name.split()[-1])
                    if 1 <= n <= 6:
                        level = n
                except Exception:
                    level = None
            if level is not None:
                lines.append("#" * level + " " + txt.strip())
            else:
                lines.append(txt)
        # tables
        for t in d.tables:
            for row in t.rows:
                cell_text = " | ".join([c.text.strip() for c in row.cells if c.text])
                if cell_text:
                    lines.append(cell_text)
        return "\n".join(lines)
    except Exception:
        return ""


def _ppt_to_text(path: str) -> str:
    # Support both .pptx (python-pptx) and legacy .ppt via fallback if possible
    suf = Path(path).suffix.lower()
    if suf == ".pptx":
        try:
            from pptx import Presentation  # type: ignore
        except Exception:
            return ""
        try:
            prs = Presentation(path)
            lines: List[str] = []
            for idx, slide in enumerate(prs.slides, start=1):
                lines.append(f"=== Slide {idx} ===")
                # slide title if available
                title_text = ""
                try:
                    if getattr(slide, "shapes", None):
                        if getattr(slide.shapes, "title", None) and slide.shapes.title and slide.shapes.title.text:
                            title_text = slide.shapes.title.text.strip()
                except Exception:
                    title_text = ""
                if title_text:
                    lines.append("# " + title_text)
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            txt = "".join(run.text for run in p.runs) or p.text
                            if txt:
                                lines.append(txt)
            return "\n".join(lines)
        except Exception:
            return ""
    else:
        # .ppt legacy: try using python-pptx won't work; if unoconv/libreoffice env not guaranteed, return empty
        return ""


def _read_file_text(path: str) -> str:
    p = Path(path)
    suf = p.suffix.lower()
    if suf in [".txt", ".md"]:
        return p.read_text(encoding="utf-8", errors="ignore")
    if suf in [".html", ".htm"]:
        return _html_to_text(p.read_text(encoding="utf-8", errors="ignore"))
    if suf == ".pdf":
        return _pdf_to_text(str(p))
    if suf == ".docx":
        return _docx_to_text(str(p))
    if suf in [".ppt", ".pptx"]:
        return _ppt_to_text(str(p))
    if suf == ".doc":
        # Try Apache Tika if available for legacy .doc
        try:
            from tika import parser  # type: ignore
            parsed = parser.from_file(str(p))
            return (parsed.get("content") or "").strip()
        except Exception:
            return ""
    return ""


def standardize_items(items: List[Dict]) -> List[Dict]:
    """
    输入 items，每项可包含：title、url、text、html、path
    输出标准文档：{title, url, content}
    - 支持解析类型：txt/md/html/htm/pdf/docx/ppt/pptx（doc 需 tika 可选）
    """
    docs: List[Dict] = []
    for it in items or []:
        title = (it.get("title") or "").strip()
        url = (it.get("url") or "").strip()
        content = (it.get("text") or "").strip()
        if not content:
            html = it.get("html")
            if html:
                content = _html_to_text(html)
        if not content:
            path = it.get("path")
            if path:
                content = _read_file_text(path)
                if not url:
                    url = f"file://{path}"
        if content:
            docs.append({"title": title, "url": url, "content": content})
    return docs